import datetime
import glob
import os
import platform
import tkinter as tk
from pptx import Presentation
from tkinter import filedialog

def clear_screen():
    if platform.system() == 'Windows':
        os.system('cls')
    else:
        os.system('clear')

def file_format_size(size):
    power = 2**10
    n = 0
    size_labels = {0: 'B', 1: 'KB', 2: 'MB', 3: 'GB', 4: 'TB'}
    while size > power:
        size /= power
        n += 1
    return f'{size:.2f} {size_labels[n]}'

def count_slides():
    folder = folder_entry.get()
    if not folder:
        folder_path = os.path.dirname(os.path.abspath(__file__))
        if folder_path:
            folder_entry.delete(0, tk.END)
            folder_entry.insert(tk.END, folder_path)
            folder = folder_path
        else:
            return
    presentations = glob.glob(os.path.join(folder, '*.ppt*'))
    total_number_of_slides = 0
    total_filesize = 0
    result_str = ""
    if presentations:
        result_str += f"Scan results of '{folder}':\n"
        for presentation in presentations:
            current_presentation = Presentation(presentation)
            number_of_slides = len(current_presentation.slides)
            file_size = os.path.getsize(presentation)
            result_str += f"Number of slides in '{os.path.basename(presentation)}': {number_of_slides} ({file_format_size(file_size)})\n"
            total_number_of_slides += number_of_slides
            total_filesize += file_size
        result_str += f"Total number of slides in {len(presentations)} presentation(s): {total_number_of_slides} ({file_format_size(total_filesize)})\n"
    else:
        result_str += "Error: There are no presentations in this folder.\n"
    result_text.delete(1.0, tk.END)
    result_text.insert(tk.END, result_str)

def browse_folder():
    folder_path = filedialog.askdirectory()
    if folder_path:
        clear_results()
        folder_entry.delete(0, tk.END)
        folder_entry.insert(tk.END, folder_path)

def clear_results():
    result_text.delete(1.0, tk.END)

def save_report():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    report_filename_prefix = "slide_counter_gui_report"
    current_datetime = datetime.datetime.now().strftime("_%Y-%m-%d_%H-%M-%S")
    report = os.path.join(script_dir, f"{report_filename_prefix}{current_datetime}.txt")
    with open(report, "w") as report_file:
        report_file.write(result_text.get(1.0, tk.END))
    report_filename = os.path.basename(report)
    result_text.insert(tk.END, f"Report saved to '{report_filename}'")

clear_screen()

root = tk.Tk()
root.geometry("960x540")
root.title("Slide Counter")
text_frame = tk.Frame(root)
text_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True)
result_text = tk.Text(text_frame)
result_text.pack(side=tk.TOP, fill=tk.BOTH, expand=True)
button_frame = tk.Frame(root)
button_frame.pack(side=tk.BOTTOM, fill=tk.X, expand=False)
clear_button = tk.Button(button_frame, text="Clear Results", command=clear_results, width=10)
clear_button.pack(side=tk.LEFT, padx=(10, 5), pady=10)
save_button = tk.Button(button_frame, text="Save Report", command=save_report, width=10)
save_button.pack(side=tk.LEFT, padx=(0, 5), pady=10)
browse_button = tk.Button(button_frame, text="Browse", command=browse_folder, width=10)
browse_button.pack(side=tk.LEFT, padx=(0, 5), pady=10)
folder_label = tk.Label(button_frame, text="Folder: ")
folder_label.pack(side=tk.LEFT, padx=(10, 0), pady=10)
folder_entry = tk.Entry(button_frame)
folder_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 10), pady=10)
count_button = tk.Button(button_frame, text="Count Slides", command=lambda: count_slides(), width=10)
count_button.pack(side=tk.RIGHT, padx=10, pady=10)
root.mainloop()